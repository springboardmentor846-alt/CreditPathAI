# train_all.py
# Run this from the folder that contains credit_risk_dataset.csv and Loan_Default.csv
from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import os, pickle, zipfile, sys, warnings
from sklearn.model_selection import train_test_split
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import roc_auc_score, roc_curve, precision_score, recall_score, f1_score, confusion_matrix, accuracy_score
from sklearn.ensemble import HistGradientBoostingClassifier
warnings.filterwarnings("ignore")

RESULT_DIR = Path("creditpath_results")
RESULT_DIR.mkdir(parents=True, exist_ok=True)

DATA1 = Path("credit_risk_dataset.csv")
DATA2 = Path("Loan_Default.csv")

if not DATA1.exists():
    raise FileNotFoundError("credit_risk_dataset.csv not found in current directory.")
if not DATA2.exists():
    raise FileNotFoundError("Loan_Default.csv not found in current directory.")

df1 = pd.read_csv(DATA1)
df2 = pd.read_csv(DATA2)

# Save snapshots
df1.head().to_csv(RESULT_DIR/"df1_head.csv", index=False)
df2.head().to_csv(RESULT_DIR/"df2_head.csv", index=False)

# Determine target column heuristically
possible_targets = ['loan_status', 'default', 'loan_default', 'target', 'loanstatus', 'loan_status.']
target_col = None
for t in possible_targets:
    if t in df1.columns:
        target_col = t
        break
if target_col is None:
    # try columns with small number of uniqs
    for col in df1.columns:
        if df1[col].nunique() <= 4 and df1[col].dtype == 'object':
            target_col = col
            break
if target_col is None:
    raise RuntimeError("Could not detect target column in credit_risk_dataset.csv. Please tell me the exact column name.")

# Map common textual values to binary target
y_raw_unique = df1[target_col].dropna().unique().tolist()
mapping = {}
for val in y_raw_unique:
    s = str(val).lower()
    if any(k in s for k in ['default','charged off','late','bad','npd']):
        mapping[val] = 1
    elif any(k in s for k in ['fully paid','paid','current','good','0','no']):
        mapping[val] = 0
    else:
        mapping[val] = 1 if s in ['1','yes','y'] else 0

df1['__target_raw__'] = df1[target_col]
df1['target'] = df1['__target_raw__'].map(mapping).fillna(0).astype(int)
pd.DataFrame({'raw_value': list(mapping.keys()), 'mapped': list(mapping.values())}).to_csv(RESULT_DIR/"target_mapping.csv", index=False)

# EDA (matplotlib only)
numeric_cols = df1.select_dtypes(include=[np.number]).columns.tolist()
numeric_cols = [c for c in numeric_cols if c not in ['target']]
for col in numeric_cols[:6]:
    plt.figure(figsize=(6,4))
    plt.hist(df1[col].dropna(), bins=30)
    plt.title(f"{col} distribution")
    plt.xlabel(col)
    plt.ylabel("count")
    plt.tight_layout()
    plt.savefig(RESULT_DIR/f"dist_{col}.png")
    plt.close()

# target distribution
plt.figure(figsize=(4,3))
vals = df1['target'].value_counts().sort_index()
plt.bar(vals.index.astype(str), vals.values)
plt.title("Target distribution (0 = good, 1 = default)")
plt.xlabel("target")
plt.ylabel("count")
plt.tight_layout()
plt.savefig(RESULT_DIR/"target_distribution.png")
plt.close()

# Correlation heatmap matrix (matplotlib)
corr = df1.select_dtypes(include=[np.number]).corr()
plt.figure(figsize=(8,6))
plt.imshow(corr, interpolation='nearest')
plt.colorbar()
plt.xticks(range(len(corr.columns)), corr.columns, rotation=90, fontsize=6)
plt.yticks(range(len(corr.columns)), corr.columns, fontsize=6)
plt.title("Correlation matrix (numeric features)")
plt.tight_layout()
plt.savefig(RESULT_DIR/"correlation_matrix.png")
plt.close()

# PREPROCESS
X = df1.drop(columns=[target_col, '__target_raw__', 'target'], errors='ignore')
y = df1['target']

X_enc = pd.get_dummies(X, drop_first=True)
imputer = SimpleImputer(strategy='median')
X_imp = pd.DataFrame(imputer.fit_transform(X_enc), columns=X_enc.columns)
X_imp.head().to_csv(RESULT_DIR/"preprocessed_head.csv", index=False)

# train/test split
X_train, X_test, y_train, y_test = train_test_split(X_imp, y, test_size=0.2, random_state=42, stratify=y)

scaler = StandardScaler()
X_train_scaled = scaler.fit_transform(X_train)
X_test_scaled = scaler.transform(X_test)

def evaluate_and_save(model, name):
    preds_proba = model.predict_proba(X_test_scaled)[:,1]
    preds = (preds_proba >= 0.5).astype(int)
    auc = roc_auc_score(y_test, preds_proba)
    prec = precision_score(y_test, preds, zero_division=0)
    rec = recall_score(y_test, preds, zero_division=0)
    f1 = f1_score(y_test, preds, zero_division=0)
    acc = accuracy_score(y_test, preds)
    cm = confusion_matrix(y_test, preds).tolist()
    metrics = {
        'model': name,
        'auc': float(auc),
        'precision': float(prec),
        'recall': float(rec),
        'f1': float(f1),
        'accuracy': float(acc),
        'confusion_matrix': cm
    }
    with open(RESULT_DIR/f"{name}.pkl", "wb") as f:
        pickle.dump(model, f)
    pd.DataFrame([metrics]).to_csv(RESULT_DIR/f"metrics_{name}.csv", index=False)
    return metrics, preds_proba

models_metrics = []

# Logistic Regression baseline
lr = LogisticRegression(max_iter=3000)
lr.fit(X_train_scaled, y_train)
metrics_lr, proba_lr = evaluate_and_save(lr, "logistic_regression")
models_metrics.append(metrics_lr)

# Try XGBoost
trained_xgb = False
try:
    import xgboost as xgb
    xgb_clf = xgb.XGBClassifier(n_estimators=300, learning_rate=0.05, max_depth=6, use_label_encoder=False, eval_metric='logloss')
    xgb_clf.fit(X_train_scaled, y_train)
    metrics_xgb, proba_xgb = evaluate_and_save(xgb_clf, "xgboost")
    models_metrics.append(metrics_xgb)
    trained_xgb = True
except Exception as e:
    print("XGBoost not available:", e)

# Try LightGBM
trained_lgb = False
try:
    import lightgbm as lgb
    lgb_clf = lgb.LGBMClassifier(n_estimators=400, learning_rate=0.03, num_leaves=31)
    lgb_clf.fit(X_train_scaled, y_train)
    metrics_lgb, proba_lgb = evaluate_and_save(lgb_clf, "lightgbm")
    models_metrics.append(metrics_lgb)
    trained_lgb = True
except Exception as e:
    print("LightGBM not available:", e)

# Fallback: HistGradientBoosting if others not present
if not (trained_xgb or trained_lgb):
    hgb = HistGradientBoostingClassifier(max_iter=300)
    hgb.fit(X_train_scaled, y_train)
    metrics_hgb, proba_hgb = evaluate_and_save(hgb, "hist_gradient_boosting")
    models_metrics.append(metrics_hgb)

pd.DataFrame(models_metrics).to_csv(RESULT_DIR/"models_comparison_metrics.csv", index=False)

# ROC plot
plt.figure(figsize=(6,5))
fpr, tpr, _ = roc_curve(y_test, proba_lr)
plt.plot(fpr, tpr, label=f"Logistic (AUC={metrics_lr['auc']:.3f})")
if 'proba_xgb' in locals():
    fpr2, tpr2, _ = roc_curve(y_test, proba_xgb); plt.plot(fpr2, tpr2, label=f"XGBoost (AUC={metrics_xgb['auc']:.3f})")
if 'proba_lgb' in locals():
    fpr3, tpr3, _ = roc_curve(y_test, proba_lgb); plt.plot(fpr3, tpr3, label=f"LightGBM (AUC={metrics_lgb['auc']:.3f})")
if 'proba_hgb' in locals():
    fpr4, tpr4, _ = roc_curve(y_test, proba_hgb); plt.plot(fpr4, tpr4, label=f"HGB (AUC={metrics_hgb['auc']:.3f})")
plt.plot([0,1],[0,1],'--'); plt.xlabel("False Positive Rate"); plt.ylabel("True Positive Rate")
plt.title("ROC Curves"); plt.legend(loc='lower right'); plt.tight_layout()
plt.savefig(RESULT_DIR/"roc_comparison.png"); plt.close()

# Save best confusion matrix
best = max(models_metrics, key=lambda x: x['auc'])
best_name = best['model']
pd.DataFrame(best['confusion_matrix']).to_csv(RESULT_DIR/f"best_confusion_matrix_{best_name}.csv", index=False)

# Build PPTX (optional; requires python-pptx)
try:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CreditPathAI — Mentor-ready Report"
    slide.placeholders[1].text = "Milestones 2–4: EDA, Baseline & Advanced Models"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Overview"
    tx = slide.shapes.placeholders[1].text_frame
    tx.text = f"credit_risk_dataset.csv: {df1.shape[0]} rows, {df1.shape[1]} cols\\nLoan_Default.csv: {df2.shape[0]} rows, {df2.shape[1]} cols"
    tx.add_paragraph().text = f"Detected target column: {target_col}. Unique raw values: {y_raw_unique}"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "EDA - Target Distribution"
    slide.shapes.add_picture(str(RESULT_DIR/"target_distribution.png"), Inches(1), Inches(1.2), width=Inches(6))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "ROC Comparison"
    slide.shapes.add_picture(str(RESULT_DIR/"roc_comparison.png"), Inches(1), Inches(1.2), width=Inches(6))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Best Model & Next Steps"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"Best model: {best_name} (AUC={best['auc']:.3f})"
    tf.add_paragraph().text = "Next steps: handle imbalance, explainability (SHAP), hyperparameter tuning, API + dashboard"

    pptx_path = RESULT_DIR/"CreditPathAI_report.pptx"
    prs.save(str(pptx_path))
    print("PPTX saved to", pptx_path)
except Exception as e:
    print("PPTX generation skipped/failed (python-pptx missing?) ->", e)

# Zip results
ZIP_PATH = Path("CreditPathAI_complete.zip")
with zipfile.ZipFile(ZIP_PATH, 'w', zipfile.ZIP_DEFLATED) as zf:
    for root, dirs, files in os.walk(RESULT_DIR):
        for f in files:
            zf.write(os.path.join(root, f), arcname=os.path.join("results", f))
    zf.write(DATA1, f"original/{DATA1.name}")
    zf.write(DATA2, f"original/{DATA2.name}")

print("Finished. Results in", RESULT_DIR, "zip:", ZIP_PATH)
